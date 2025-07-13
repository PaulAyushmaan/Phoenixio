from pptx import Presentation

class TopicExtraction:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.presentation = Presentation(ppt_path)
        self.segments = []

    def extract_slide_text(self):
        prs = Presentation(self.ppt_path)
        slides = []
        for idx, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
            slides.append({
                "slide_num": idx + 1,
                "text": " ".join(text)
            })
        return slides
    @staticmethod
    def get_window_size(slide_count, alpha=0.15):
        """
        Calculate context-aware window size based on slide count,
        using min-max ranges instead of hard-coded values.

        Parameters:
        - slide_count: total number of slides in the deck
        - alpha: scaling factor within the allowed range

        Returns:
        - window_size: integer window size adapted to slide count
        """
        if slide_count <= 15:
            window_size = 2
        elif slide_count <= 30:
            window_size = 3
        elif slide_count <= 50:
            window_size = 4
        else:
            window_size = 5

        return window_size
    
    def group_slides_by_window(self, window_size):
        """
        Groups slides into windows based on window size.

        Parameters:
        - slide_blocks: List of dicts with 'slide_num' and 'text' keys.
        - window_size: Number of slides per group.

        Returns:
        - List of grouped windows, each containing a list of slide dicts.
        """
        grouped_windows = []
        temp_window = []

        for block in self.extract_slide_text():
            temp_window.append(block)
            if len(temp_window) == window_size:
                grouped_windows.append(temp_window)
                temp_window = []

        if temp_window:
            grouped_windows.append(temp_window)

        return grouped_windows
    




